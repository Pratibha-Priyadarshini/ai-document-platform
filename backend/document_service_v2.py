from docx import Document
from docx.shared import Pt, Inches as DocxInches
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
from typing import List, Dict, Any
from models import Project, Section
from image_service import image_service
from theme_service import theme_service

class DocumentServiceV2:
    """Improved PowerPoint generation with better formatting and layout."""
    
    @staticmethod
    def create_professional_pptx(project: Project, sections: List[Section], metadata: Dict[str, Any] = None) -> BytesIO:
        """Create a professional PowerPoint presentation using built-in layouts or custom themes."""
        # Extract theme preference
        theme = metadata.get('theme', {}) if metadata else {}
        theme_id = theme.get('id', None)
        
        # Try to load custom theme template
        if theme_id:
            prs = theme_service.load_theme_template(theme_id)
            if prs:
                print(f"✓ Using custom theme: {theme_id}")
            else:
                print(f"ℹ Theme '{theme_id}' not found, using default with colors")
                prs = Presentation()
        else:
            prs = Presentation()
        
        # Set standard dimensions if not already set
        if not prs.slide_width or prs.slide_width == 0:
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # Available layouts:
        # 0: Title Slide
        # 1: Title and Content
        # 2: Section Header
        # 3: Two Content
        # 4: Comparison
        # 5: Title Only
        # 6: Blank
        # 7: Content with Caption
        # 8: Picture with Caption
        
        # Extract styling
        theme = metadata.get('theme', {}) if metadata else {}
        text_style = metadata.get('textStyle', {}) if metadata else {}
        
        primary_color = theme.get('preview', {}).get('primary', '#667eea')
        secondary_color = theme.get('preview', {}).get('secondary', '#764ba2')
        text_color = theme.get('preview', {}).get('text', '#ffffff')
        
        font_family = text_style.get('fontFamily', 'arial')
        alignment = text_style.get('alignment', 'left')
        image_alignment = text_style.get('imageAlignment', 'top')
        
        # Font mapping
        font_map = {
            'arial': 'Arial',
            'helvetica': 'Helvetica',
            'georgia': 'Georgia',
            'times': 'Times New Roman',
            'courier': 'Courier New',
            'verdana': 'Verdana'
        }
        font_name = font_map.get(font_family, 'Arial')
        
        # Alignment mapping
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT
        }
        text_align = align_map.get(alignment, PP_ALIGN.LEFT)
        
        # Create title slide
        DocumentServiceV2._create_title_slide(
            prs, project, primary_color, secondary_color, text_color, font_name
        )
        
        # Create content slides
        for section in sorted(sections, key=lambda x: x.order):
            DocumentServiceV2._create_content_slide(
                prs, section, project.main_topic,
                primary_color, secondary_color, text_color,
                font_name, text_align, image_alignment
            )
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream
    
    @staticmethod
    def _create_title_slide(prs, project, primary_color, secondary_color, text_color, font_name):
        """Create an attractive title slide using built-in layout."""
        # Use Title Slide layout (layout 0)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Apply gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(primary_color))
        fill.gradient_stops[1].color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(secondary_color))
        
        # Use built-in title placeholder
        title_shape = slide.shapes.title
        title_shape.text = project.title
        
        # Format title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = PptxPt(54)
        title_para.font.bold = True
        title_para.font.name = font_name
        title_para.font.color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(text_color))
        title_para.alignment = PP_ALIGN.CENTER
        
        # Always add a relevant subtitle
        subtitle_text = project.main_topic
        if not subtitle_text or subtitle_text.lower().strip() == project.title.lower().strip():
            # Generate a relevant subheading
            subtitle_text = f"A Comprehensive Overview"
        
        # Find subtitle placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = subtitle_text
                subtitle_frame = shape.text_frame
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = PptxPt(24)
                subtitle_para.font.name = font_name
                subtitle_para.font.color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(text_color))
                subtitle_para.alignment = PP_ALIGN.CENTER
                break
    
    @staticmethod
    def _create_content_slide(prs, section, main_topic, primary_color, secondary_color, 
                             text_color, font_name, text_align, image_alignment):
        """Create a content slide using built-in Title and Content layout."""
        # Use Title and Content layout (layout 1) or Picture with Caption (layout 8)
        if image_alignment in ['left', 'right']:
            slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content layout
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        
        # Apply gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        fill.gradient_stops[0].color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(primary_color))
        fill.gradient_stops[1].color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(secondary_color))
        
        # Use built-in title placeholder
        title_shape = slide.shapes.title
        title_shape.text = section.title
        
        # Format title
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.size = PptxPt(36)
        title_para.font.bold = True
        title_para.font.name = font_name
        title_para.font.color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(text_color))
        title_para.alignment = PP_ALIGN.LEFT
        
        # Find content placeholder by index
        content_placeholder = None
        for shape in slide.placeholders:
            # Content placeholder is usually index 1 (after title which is 0)
            if shape.placeholder_format.idx == 1:
                content_placeholder = shape
                break
        
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear default text
            text_frame.word_wrap = True
            
            # Clean and add content with strict limits
            # Fewer lines for side-by-side layouts
            max_lines = 5 if image_alignment in ['left', 'right'] else 6
            max_words = 16 if image_alignment in ['left', 'right'] else 18
            
            lines = DocumentServiceV2._clean_content(section.content, max_lines, max_words, section.title)
            
            if lines:
                # Adjust content placeholder position for top/bottom layouts
                if image_alignment == 'top':
                    # Move content down to make room for image at top
                    # Image will be placed above, so content stays in default position
                    pass
                elif image_alignment == 'bottom':
                    # Content stays at top, image will be placed below
                    pass
                
                for i, line in enumerate(lines):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = line
                    p.font.size = PptxPt(16)  # Slightly smaller font
                    p.font.name = font_name
                    p.font.color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(text_color))
                    p.alignment = text_align
                    p.space_before = PptxPt(4)  # Reduced spacing
                    p.space_after = PptxPt(4)
                    p.level = 0
            else:
                # If no content, add a message
                p = text_frame.paragraphs[0]
                p.text = "No content available"
                p.font.size = PptxPt(18)
                p.font.name = font_name
                p.font.color.rgb = RGBColor(*DocumentServiceV2._hex_to_rgb(text_color))
        
        # Calculate text content height for dynamic image sizing
        text_height = 0
        num_lines = 0
        if content_placeholder:
            # Estimate text height based on number of lines and font size
            num_lines = len(lines) if lines else 0
            line_height = PptxPt(16 + 8)  # font size + spacing
            text_height = num_lines * line_height
        
        # Add image for ALL layouts with dynamic sizing
        try:
            print(f"\n🖼️  Fetching image for slide: '{section.title}' (text lines: {num_lines})")
            image_data = image_service.generate_slide_image(section.title, main_topic)
            
            if image_data:
                print(f"✓ Image data received, calculating optimal size...")
                
                if image_alignment in ['left', 'right']:
                    # Find second content placeholder for image
                    for shape in slide.placeholders:
                        if shape.placeholder_format.idx == 2:  # Second content area
                            # Use placeholder dimensions (already optimized for side layout)
                            left = shape.left
                            top = shape.top
                            width = shape.width
                            height = shape.height
                            
                            # Remove placeholder and add image
                            sp = shape.element
                            sp.getparent().remove(sp)
                            
                            slide.shapes.add_picture(image_data, left, top, width=width, height=height)
                            print(f"✓ Image added to {image_alignment} layout (size: {width/914400:.1f}x{height/914400:.1f} inches)")
                            break
                else:
                    # For top/bottom layouts, calculate optimal dimensions based on text
                    dims = DocumentServiceV2._calculate_image_dimensions(prs, num_lines, image_alignment)
                    
                    slide.shapes.add_picture(
                        image_data, 
                        dims['left'], 
                        dims['top'], 
                        width=dims['width'], 
                        height=dims['height']
                    )
                    print(f"✓ Image added to {image_alignment} layout (size: {dims['width']/914400:.1f}x{dims['height']/914400:.1f} inches, top: {dims['top']/914400:.1f})")
            else:
                print(f"✗ No image data received for '{section.title}'")
        except Exception as e:
            print(f"✗ Image error for '{section.title}': {e}")
    
    @staticmethod
    def _calculate_image_dimensions(prs, num_lines, image_alignment):
        """Calculate optimal image dimensions based on text content."""
        from pptx.util import Inches
        
        slide_height = prs.slide_height
        title_height = Inches(1.2)
        
        # Calculate text space needed (approximately 0.35 inches per line)
        text_space = Inches(0.5 + (num_lines * 0.35))
        
        if image_alignment == 'top':
            # Image at top, text below
            available_height = slide_height - title_height - text_space - Inches(0.5)
            top = title_height
            height = max(Inches(1.5), min(available_height, Inches(3.5)))
        else:  # bottom
            # Text at top, image below
            top = title_height + text_space + Inches(0.3)
            available_height = slide_height - top - Inches(0.4)
            height = max(Inches(1.5), min(available_height, Inches(3)))
        
        return {
            'left': Inches(1.5),
            'top': top,
            'width': Inches(7),
            'height': height
        }
    
    @staticmethod
    def _get_layout_config(image_alignment):
        """Get layout configuration for different image alignments (legacy - for reference)."""
        layouts = {
            'top': {
                'image': {'left': 1.5, 'top': 1.4, 'width': 7, 'height': 2.8},
                'content': {'left': 0.75, 'top': 4.4, 'width': 8.5, 'height': 2.9}
            },
            'bottom': {
                'content': {'left': 0.75, 'top': 1.4, 'width': 8.5, 'height': 2.9},
                'image': {'left': 1.5, 'top': 4.5, 'width': 7, 'height': 2.8}
            },
            'left': {
                'image': {'left': 0.3, 'top': 1.4, 'width': 4.2, 'height': 5.8},
                'content': {'left': 4.7, 'top': 1.4, 'width': 5, 'height': 5.9}
            },
            'right': {
                'content': {'left': 0.5, 'top': 1.4, 'width': 5, 'height': 5.9},
                'image': {'left': 5.7, 'top': 1.4, 'width': 4.2, 'height': 5.8}
            }
        }
        return layouts.get(image_alignment, layouts['top'])
    
    @staticmethod
    def _clean_content(content, max_lines=5, max_words_per_line=18, section_title=None):
        """Clean and format content lines with strict limits to prevent overflow (for PowerPoint)."""
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        
        cleaned = []
        for line in lines:
            # Skip markers
            if line in ['**', '*', '•', '-', '---', '+', '']:
                continue
            
            # Remove markdown and bullets
            line = line.replace('**', '').replace('*', '')
            while line and line[0] in ['•', '-', '*', '·', '+', ' ']:
                line = line[1:].strip()
            
            # Skip if line matches section title (case-insensitive, partial match)
            if section_title:
                line_lower = line.lower().strip()
                title_lower = section_title.lower().strip()
                # Exact match or line contains the title
                if (line_lower == title_lower or 
                    line_lower.startswith(title_lower) or
                    line_lower.endswith(title_lower)):
                    continue
            
            # Skip headers and empty lines
            if (line.lower().startswith('slide title:') or 
                line.lower().startswith('*slide title') or
                line.lower().startswith('section title:') or
                line.lower().startswith('title:') or
                line.lower().startswith('here are') or
                line.lower().startswith('here is') or
                line.lower().startswith('introduction to') or
                (line.lower().startswith('key ') and line.endswith(':**')) or
                len(line) < 3):
                continue
            
            # Limit line length - split if too long
            words = line.split()
            if len(words) > max_words_per_line:
                # Split into multiple lines
                current_line = []
                for word in words:
                    current_line.append(word)
                    if len(current_line) >= max_words_per_line:
                        cleaned.append(' '.join(current_line))
                        current_line = []
                        if len(cleaned) >= max_lines:
                            break
                if current_line and len(cleaned) < max_lines:
                    cleaned.append(' '.join(current_line))
            else:
                cleaned.append(line)
            
            # Stop if we have enough lines
            if len(cleaned) >= max_lines:
                break
        
        return cleaned[:max_lines]
    
    @staticmethod
    def _clean_content_for_word(content, section_title=None):
        """Clean content for Word documents - no line limits, preserve full content."""
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        
        cleaned = []
        for line in lines:
            # Skip markers
            if line in ['**', '*', '•', '-', '---', '+', '']:
                continue
            
            # Remove markdown formatting
            line = line.replace('**', '').replace('*', '')
            
            # Skip if line matches section title (case-insensitive, partial match)
            if section_title:
                line_lower = line.lower().strip()
                title_lower = section_title.lower().strip()
                if (line_lower == title_lower or 
                    line_lower.startswith(title_lower) or
                    line_lower.endswith(title_lower)):
                    continue
            
            # Skip headers and empty lines
            if (line.lower().startswith('slide title:') or 
                line.lower().startswith('section title:') or
                line.lower().startswith('title:') or
                line.lower().startswith('here are') or
                line.lower().startswith('here is') or
                len(line) < 3):
                continue
            
            cleaned.append(line)
        
        return cleaned
    
    @staticmethod
    def _add_image_to_doc(doc, image_data, width, height):
        """Add an image to the document with specified dimensions."""
        try:
            para = doc.add_paragraph()
            para.alignment = 1  # Center alignment
            run = para.add_run()
            run.add_picture(image_data, width=width, height=height)
            print(f"✓ Image added to document (width: {width/914400:.1f}\", height: {height/914400:.1f}\")")
        except Exception as e:
            print(f"✗ Failed to add image to document: {e}")
    
    @staticmethod
    def _hex_to_rgb(hex_color):
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    @staticmethod
    def _apply_document_theme(doc, font_name, primary_rgb, secondary_rgb, accent_rgb):
        """Apply theme colors to document styles."""
        from docx.shared import RGBColor as DocxRGBColor
        
        try:
            # Update Heading 1 style
            if 'Heading 1' in doc.styles:
                heading1_style = doc.styles['Heading 1']
                heading1_style.font.color.rgb = DocxRGBColor(*accent_rgb)
                heading1_style.font.name = font_name
                heading1_style.font.size = Pt(18)
                heading1_style.font.bold = True
            
            # Update Title style
            if 'Title' in doc.styles:
                title_style = doc.styles['Title']
                title_style.font.color.rgb = DocxRGBColor(*primary_rgb)
                title_style.font.name = font_name
                title_style.font.size = Pt(28)
                title_style.font.bold = True
            
            # Update Normal style
            if 'Normal' in doc.styles:
                normal_style = doc.styles['Normal']
                normal_style.font.name = font_name
                normal_style.font.size = Pt(12)
        except Exception as e:
            print(f"Warning: Could not apply document theme: {e}")

    @staticmethod
    def create_docx(project: Project, sections: List[Section], metadata: Dict[str, Any] = None) -> BytesIO:
        """Create a Word document from project data with theme and style support."""
        from docx.shared import RGBColor as DocxRGBColor
        
        # Extract theme and text style preferences
        theme = metadata.get('theme', {}) if metadata else {}
        text_style = metadata.get('textStyle', {}) if metadata else {}
        theme_id = theme.get('id', None)
        theme_preview = theme.get('preview', {}) if theme else {}
        
        # Get theme colors
        primary_color = theme_preview.get('primary', '#0078D4')
        secondary_color = theme_preview.get('secondary', '#106EBE')
        accent_color = theme_preview.get('accent', '#50E6FF')
        
        # Convert hex to RGB
        primary_rgb = DocumentServiceV2._hex_to_rgb(primary_color)
        secondary_rgb = DocumentServiceV2._hex_to_rgb(secondary_color)
        accent_rgb = DocumentServiceV2._hex_to_rgb(accent_color)
        
        # Text style settings
        font_name = text_style.get('fontFamily', 'calibri').title()
        if font_name == 'Times': font_name = 'Times New Roman'
        elif font_name == 'Courier': font_name = 'Courier New'
        
        font_size = text_style.get('fontSize', 12)
        alignment_map = {'left': 0, 'center': 1, 'right': 2, 'justify': 3}
        text_alignment = alignment_map.get(text_style.get('alignment', 'left'), 0)
        image_alignment = text_style.get('imageAlignment', 'top')
        
        # Try to load custom theme template
        if theme_id:
            doc = theme_service.load_theme_template(theme_id, 'docx')
            if doc:
                print(f"✓ Using Word theme: {theme_id}")
            else:
                print(f"ℹ Theme '{theme_id}' not found, using default")
                doc = Document()
        else:
            doc = Document()
        
        # Apply document-level theme
        DocumentServiceV2._apply_document_theme(doc, font_name, primary_rgb, secondary_rgb, accent_rgb)
        
        # Create styled header section with theme colors
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        # Add title with theme background
        title_para = doc.add_paragraph()
        title_para.alignment = 1  # Center alignment
        
        # Add shading (background color) to title paragraph
        try:
            shading_elm = OxmlElement('w:shd')
            shading_elm.set(qn('w:fill'), primary_color.lstrip('#'))
            title_para._element.get_or_add_pPr().append(shading_elm)
        except:
            pass
        
        # Add title text
        title_run = title_para.add_run(project.title)
        title_run.font.color.rgb = DocxRGBColor(255, 255, 255)  # White text on colored background
        title_run.font.bold = True
        title_run.font.size = Pt(28)
        title_run.font.name = font_name
        
        # Add spacing before and after
        title_para.paragraph_format.space_before = Pt(12)
        title_para.paragraph_format.space_after = Pt(12)
        

        

        
        # Add a horizontal line separator
        separator = doc.add_paragraph()
        separator_run = separator.add_run('_' * 80)
        separator_run.font.color.rgb = DocxRGBColor(*primary_rgb)
        separator.alignment = 1
        
        doc.add_paragraph()  # Empty line
        
        # Add sections with custom styling and images
        for section in sorted(sections, key=lambda x: x.order):
            # Section heading with theme color and background
            heading = doc.add_paragraph()
            
            # Add background color to section heading
            try:
                from docx.oxml.ns import qn
                from docx.oxml import OxmlElement
                shading_elm = OxmlElement('w:shd')
                shading_elm.set(qn('w:fill'), accent_color.lstrip('#'))
                heading._element.get_or_add_pPr().append(shading_elm)
            except:
                pass
            
            heading_run = heading.add_run(section.title)
            heading_run.font.color.rgb = DocxRGBColor(255, 255, 255)  # White text
            heading_run.font.bold = True
            heading_run.font.size = Pt(18)
            heading_run.font.name = font_name
            
            heading.paragraph_format.space_before = Pt(12)
            heading.paragraph_format.space_after = Pt(6)
            
            # Fetch image for this section
            print(f"\n🖼️  Fetching image for section: '{section.title}'")
            image_data = None
            try:
                image_data = image_service.generate_slide_image(section.title, project.main_topic)
                if image_data:
                    print(f"✓ Image data received for '{section.title}'")
            except Exception as e:
                print(f"✗ Image fetch failed for '{section.title}': {e}")
            
            # Clean content for Word documents (no line limits)
            lines = DocumentServiceV2._clean_content_for_word(section.content, section_title=section.title)
            
            # Add content and image based on alignment
            if image_alignment == 'top' and image_data:
                # Image at top, then content
                DocumentServiceV2._add_image_to_doc(doc, image_data, DocxInches(5), DocxInches(3))
                doc.add_paragraph()  # Space after image
                
                for line in lines:
                    para = doc.add_paragraph(line)
                    para.alignment = text_alignment
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)
            
            elif image_alignment == 'bottom' and image_data:
                # Content first, then image
                for line in lines:
                    para = doc.add_paragraph(line)
                    para.alignment = text_alignment
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)
                
                doc.add_paragraph()  # Space before image
                DocumentServiceV2._add_image_to_doc(doc, image_data, DocxInches(5), DocxInches(3))
            
            elif image_alignment in ['left', 'right'] and image_data:
                # Side-by-side layout using table
                table = doc.add_table(rows=1, cols=2)
                table.autofit = False
                table.allow_autofit = False
                
                if image_alignment == 'left':
                    # Image in left cell, text in right cell
                    img_cell = table.rows[0].cells[0]
                    text_cell = table.rows[0].cells[1]
                    img_cell.width = DocxInches(3)
                    text_cell.width = DocxInches(4)
                else:
                    # Text in left cell, image in right cell
                    text_cell = table.rows[0].cells[0]
                    img_cell = table.rows[0].cells[1]
                    text_cell.width = DocxInches(4)
                    img_cell.width = DocxInches(3)
                
                # Add image to image cell
                img_para = img_cell.paragraphs[0]
                run = img_para.add_run()
                run.add_picture(image_data, width=DocxInches(2.8))
                
                # Add text to text cell
                text_cell.text = ''  # Clear default text
                for line in lines:
                    para = text_cell.add_paragraph(line)
                    para.alignment = text_alignment
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)
            
            else:
                # No image or image failed - just add text
                for line in lines:
                    para = doc.add_paragraph(line)
                    para.alignment = text_alignment
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size)
            
            doc.add_paragraph()  # Empty line between sections
        
        # Save to BytesIO
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream
